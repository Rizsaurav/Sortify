"""
Upload API routes - Handles document uploads and task management.
Clean separation of concerns with proper dependency injection.
"""

import io
import uuid
import asyncio
from datetime import datetime
from typing import Optional

from fastapi import APIRouter, HTTPException, UploadFile, File, Form, BackgroundTasks, Query

from models import DocumentUploadResponse, TaskStatusResponse, FileCategoryResponse
from services import get_document_service, get_categorization_service
from core import get_database_service
from utils import get_logger
from api.task_manager import get_task_manager

logger = get_logger(__name__)

# Create router
router = APIRouter(
    prefix="/upload",
    tags=["uploads"],
    responses={404: {"description": "Not found"}},
)


@router.post("", response_model=DocumentUploadResponse)
async def upload_document(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    user_id: str = Form(...)
):
    """
    Upload a document and queue it for processing.
    Returns immediately without waiting for ML processing.
    
    Args:
        file: Uploaded file
        user_id: User ID
    
    Returns:
        Upload response with task ID
    """
    try:
        logger.info(f"Upload request: {file.filename} (user={user_id})")
        
        # Get services
        doc_service = get_document_service()
        db_service = get_database_service()
        task_manager = get_task_manager()
        
        # Read file content
        content_bytes = await file.read()
        
        # Extract text based on file type
        content_str = await _extract_content(file, content_bytes)

        #Get unique filename to avoid filename conflicts

        unique_filename = _get_unique_filename(file.filename, user_id, db_service)

        #Log if filename was changed
        if unique_filename != file.filename:
            logger.info(f"Renamed '{file.filename}' to '{unique_filename}' to avoid conflict")
        
        # Check for duplicates
        duplicate_id = doc_service.check_duplicate(content_str, user_id)
        if duplicate_id:
            logger.info(f"Duplicate detected: {duplicate_id}")
            return DocumentUploadResponse(
                filename=unique_filename,
                status="duplicate",
                message=f"Document already exists with ID {duplicate_id}",
                doc_id=duplicate_id,
                task_id=None,
                timestamp=datetime.now()
            )

        # Construct storage path for original file (user_id/filename)
        storage_path = f"{user_id}/{unique_filename}"

        # Upload original binary file to Supabase Storage (best-effort).
        # IMPORTANT: frontend expects bucket 'user-files' and path 'user_id/filename'.
        file_url: Optional[str] = None
        try:
            file_url = db_service.upload_file_to_bucket(
                bucket="user-files",
                path=storage_path,
                data=content_bytes,
                content_type=file.content_type or "application/octet-stream",
            )
        except Exception as e:
            # Do not fail the upload if storage upload fails; just log
            logger.error(f"Failed to upload original file to storage: {e}")

        # Insert parent document to database (with FULL content)
        doc_id = db_service.insert_document(
            content=content_str,  # FULL CONTENT, not preview!
            user_id=user_id,  # Top-level user_id for efficient querying
            metadata={
                'user_id': user_id,  # Also keep in metadata for backward compatibility
                'filename': unique_filename,
                'original_filename': file.filename,
                'type': file.content_type,
                'size': len(content_bytes)
            },
            embedding=None,  # Will be set after chunking
            cluster_id=None,  # Will be set after categorization
            storage_path=storage_path,
            file_path=storage_path,  # Same as storage_path for now
            file_url=file_url  # Public URL for original file (if available)
        )
        
        logger.info(f"Created document {doc_id} with {len(content_str)} characters")
        
        # Generate task ID
        task_id = str(uuid.uuid4())
        
        # Add to task queue for chunking + categorization
        task_manager.add_task(
            task_id=task_id,
            doc_id=doc_id,
            user_id=user_id,
            content=content_str,  # Full content for processing
            filename=unique_filename,
            file_type=file.content_type or 'unknown',
            file_size=len(content_bytes)
        )
        
        # Kick off async processing (chunking + embeddings + categorization)
        asyncio.create_task(task_manager.process_queue())
        
        logger.info(f"Task {task_id} queued for processing")
        
        return DocumentUploadResponse(
            filename=unique_filename,
            status="queued",
            message=f"Document uploaded. Task {task_id} queued for processing.",
            doc_id=doc_id,
            task_id=task_id,
            timestamp=datetime.now()
        )
    
    except Exception as e:
        logger.error(f"Upload failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")


@router.get("/documents")
async def get_documents(user_id: str = Query(...)):
    """
    Get all documents for a user.
    
    Args:
        user_id: User ID
    
    Returns:
        List of documents
    """
    try:
        db_service = get_database_service()
        documents = db_service.get_documents_by_user(user_id)
        
        return {
            "success": True,
            "documents": documents,
            "count": len(documents)
        }
    except Exception as e:
        logger.error(f"Failed to get documents: {e}")
        raise HTTPException(status_code=500, detail="Failed to fetch documents")


@router.get("/task-status/{task_id}", response_model=TaskStatusResponse)
async def get_task_status(task_id: str):
    """
    Get the status of a processing task.
    
    Args:
        task_id: Task ID
    
    Returns:
        Task status information
    """
    try:
        task_manager = get_task_manager()
        task = task_manager.get_task(task_id)
        
        if not task:
            raise HTTPException(status_code=404, detail="Task not found")
        
        # Get category name if available
        category_name = None
        if task.category_id:
            db_service = get_database_service()
            categories = db_service.get_categories_by_user(task.user_id)
            for cat in categories:
                if cat['id'] == task.category_id:
                    category_name = cat['label']
                    break
        
        return TaskStatusResponse(
            task_id=task.task_id,
            doc_id=task.doc_id,
            status=task.status.value,
            created_at=task.created_at.isoformat(),
            category_id=task.category_id,
            category_name=category_name,
            completed_at=task.completed_at.isoformat() if task.completed_at else None,
            error=task.error
        )
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to get task status: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/file-category/{doc_id}", response_model=FileCategoryResponse)
async def get_file_category(doc_id: str):
    """
    Get the category of a file.
    
    Args:
        doc_id: Document ID
    
    Returns:
        File category information
    """
    try:
        db_service = get_database_service()
        
        # Get document
        doc = db_service.get_document(doc_id)
        if not doc:
            raise HTTPException(status_code=404, detail="Document not found")
        
        # Get category if assigned
        cluster_id = doc.get('cluster_id')
        category_name = "Uncategorized"
        
        if cluster_id:
            user_id = doc.get('metadata', {}).get('user_id')
            categories = db_service.get_categories_by_user(user_id)
            for cat in categories:
                if cat['id'] == cluster_id:
                    category_name = cat['label']
                    break
        
        filename = doc.get('metadata', {}).get('filename', 'Unknown')
        
        return FileCategoryResponse(
            doc_id=doc_id,
            category=category_name,
            cluster_id=cluster_id,
            filename=filename,
            status="categorized" if cluster_id else "pending"
        )
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to get file category: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# Helper function
async def _extract_content(file: UploadFile, content_bytes: bytes) -> str:
    """Extract text content from file based on type."""
    try:
        if file.content_type and file.content_type.startswith('text/'):
            return content_bytes.decode("utf-8", errors="ignore")
        
        elif file.content_type == 'application/pdf':
            try:
                from pypdf import PdfReader
                pdf_file = io.BytesIO(content_bytes)
                reader = PdfReader(pdf_file)
                text_pages = [
                    page.extract_text()
                    for page in reader.pages
                    if page.extract_text()
                ]
                content = "\n\n".join(text_pages)
                return content if content.strip() else f"PDF: {file.filename} (no text)"
            except Exception as e:
                logger.error(f"PDF extraction failed: {e}")
                return f"PDF: {file.filename} (extraction failed)"
        
        elif file.content_type and file.content_type.startswith('image/'):
            return f"Image: {file.filename}\nType: {file.content_type}\nSize: {len(content_bytes)} bytes"
        
        # Basic PPTX support (best-effort; falls back on failure)
        elif (
            file.content_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ) or (file.filename and file.filename.lower().endswith(".pptx")):
            try:
                from pptx import Presentation  # type: ignore

                prs = Presentation(io.BytesIO(content_bytes))
                slide_texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        text = getattr(shape, "text", None)
                        if text and text.strip():
                            slide_texts.append(text.strip())

                content = "\n\n".join(slide_texts)
                return content if content.strip() else f"PPTX: {file.filename} (no text)"
            except Exception as e:
                logger.error(f"PPTX extraction failed: {e}")
                return f"PPTX: {file.filename} (extraction failed)"
        
        else:
            return f"File: {file.filename}\nType: {file.content_type or 'unknown'}"
    
    except Exception as e:
        logger.error(f"Content extraction failed: {e}")
        return f"File: {file.filename} (extraction failed)"
    

def _get_unique_filename(filename: str, user_id: str, db_service) -> str:
    """
    Generate a unique filename by appending numbers if needed.
    Example: document.pdf -> document (1).pdf -> document (2).pdf
    
    Args:
        filename: Original filename
        user_id: User ID
        db_service: Database service instance
    
    Returns:
        Unique filename
    """
    import os
    
    # Get all existing filenames for this user
    existing_docs = db_service.get_documents_by_user(user_id)
    existing_filenames = {
        doc.get('metadata', {}).get('filename', '') 
        for doc in existing_docs
    }
    
    # If filename doesn't exist, return original
    if filename not in existing_filenames:
        return filename
    
    # Split filename into name and extension
    name, ext = os.path.splitext(filename)
    
    # Find next available number
    counter = 1
    new_filename = f"{name} ({counter}){ext}"
    
    while new_filename in existing_filenames:
        counter += 1
        new_filename = f"{name} ({counter}){ext}"
    
    return new_filename



@router.post("/initialize-categories")
async def initialize_user_categories(user_id: str):
    """
    Initialize standard categories for a user.
    Useful for existing users or manual category setup.
    
    Args:
        user_id: User ID
    
    Returns:
        List of created category IDs and names
    """
    try:
        from services import get_improved_categorization_service
        
        cat_service = get_improved_categorization_service()
        category_ids = cat_service.initialize_standard_categories(user_id)
        
        if category_ids:
            return {
                "success": True,
                "message": f"Initialized {len(category_ids)} standard categories",
                "category_count": len(category_ids),
                "categories": cat_service.standard_categories
            }
        else:
            return {
                "success": False,
                "message": "Categories already exist or initialization failed"
            }
    
    except Exception as e:
        logger.error(f"Category initialization failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))




